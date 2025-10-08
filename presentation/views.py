
from django.http import HttpResponse
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import os
from PIL import Image

def download_presentation(request):
    prs = Presentation()

    # Slide 1 - Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Civil Safety and Emergency Response Plans"
    slide.placeholders[1].text = "Educational Lecture Series\n2025"

    # Slide 2 - Objectives
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Objectives"
    content = slide.placeholders[1]
    content.text = (
        "- Understand Ukrainian civil safety legislation\n"
        "- Learn PPLAS, PLLA, PLA planning\n"
        "- Explore emergency management procedures"
    )

    # Slides with pictures
    image_files = ["lecture1.jpg", "lecture2.jpg", "lecture3.jpg"]
    for idx, img in enumerate(image_files, start=3):
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Emergency Plan Example {idx-2}"
        image_path = os.path.join('static', 'images', img)

        if os.path.exists(image_path):
            # Перевірка та конвертація WEBP у JPEG
            try:
                with Image.open(image_path) as im:
                    if im.format not in ['BMP', 'GIF', 'JPEG', 'PNG', 'TIFF', 'WMF']:
                        temp_path = os.path.join('static', 'images', img.split('.')[0] + "_converted.jpg")
                        im.convert("RGB").save(temp_path, "JPEG")
                        image_path = temp_path
            except Exception as e:
                tx_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
                tf = tx_box.text_frame
                tf.text = f"(Error opening {img}: {str(e)})"
                continue

            # Додаємо картинку у слайд
            slide.shapes.add_picture(image_path, Inches(1), Inches(2), width=Inches(6))
        else:
            # fallback, якщо файл не знайдено
            tx_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
            tf = tx_box.text_frame
            tf.text = f"(Image {img} not found)"

    # Slides 6–30 - placeholder content
    for i in range(6, 31):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Slide {i}"
        slide.placeholders[1].text = f"Content for slide {i} goes here."

    # Save presentation in memory
    output = BytesIO()
    prs.save(output)
    output.seek(0)

    # Return as downloadable file
    response = HttpResponse(
        output.read(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    response['Content-Disposition'] = 'attachment; filename="Emergency_Response_Plans.pptx"'
    return response
